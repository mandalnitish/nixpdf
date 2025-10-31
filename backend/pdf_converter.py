import sys
from pdf2docx import Converter
from PyPDF2 import PdfReader
from docx import Document

# Ensure UTF-8 output
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def pdf_to_word(pdf_path, word_path):
    try:
        print("[INFO] Starting PDF to Word conversion...")
        cv = Converter(pdf_path)
        cv.convert(word_path, start=0, end=None)
        cv.close()
        print("[SUCCESS] Conversion completed successfully.")
    except Exception as e:
        print(f"[ERROR] with pdf2docx: {e}")
        print("[INFO] Trying fallback method (simple text extraction)...")

        try:
            reader = PdfReader(pdf_path)
            doc = Document()
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    doc.add_paragraph(text)
            doc.save(word_path)
            print("[SUCCESS] Fallback method completed successfully.")
        except Exception as e2:
            print(f"[ERROR] (simple method): {e2}")
            raise e2

def pdf_to_pptx(pdf_path, pptx_path):
    from pptx import Presentation
    prs = Presentation()
    reader = PdfReader(pdf_path)
    for page in reader.pages:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        text = page.extract_text() or ""
        slide.shapes.title.text = text[:100] if len(text) > 100 else "Page Text"
        textbox = slide.shapes.add_textbox(0, 100000, 9144000, 6000000)
        textbox.text = text
    prs.save(pptx_path)
    print("[SUCCESS] PDF to PowerPoint conversion complete.")

def pdf_to_excel(pdf_path, excel_path):
    import pandas as pd
    reader = PdfReader(pdf_path)
    data = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        data.append({"Page": i + 1, "Content": text})
    df = pd.DataFrame(data)
    df.to_excel(excel_path, index=False)
    print("[SUCCESS] PDF to Excel conversion complete.")

if __name__ == "__main__":
    if len(sys.argv) != 4:
        print("Usage: python pdf_converter.py <format> <input.pdf> <output>")
        sys.exit(1)

    fmt, input_file, output_file = sys.argv[1], sys.argv[2], sys.argv[3]

    if fmt == "word":
        pdf_to_word(input_file, output_file)
    elif fmt == "pptx":
        pdf_to_pptx(input_file, output_file)
    elif fmt == "excel":
        pdf_to_excel(input_file, output_file)
    else:
        print("Unsupported format. Use: word | pptx | excel")
        sys.exit(1)
